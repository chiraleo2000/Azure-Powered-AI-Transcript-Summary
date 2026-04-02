import os
import json
import tempfile
import subprocess
from typing import Optional, Dict, Any, List
import PyPDF2
import docx
from openpyxl import load_workbook
import pandas as pd
from pptx import Presentation
from PIL import Image
import zipfile
import csv
import chardet

class FileProcessor:
    """Enhanced file processor for various document types with improved error handling"""
    
    def __init__(self):
        self.supported_extensions = {
            'pdf': self._process_pdf,
            'docx': self._process_docx,
            'doc': self._process_doc,
            'pptx': self._process_pptx,
            'ppt': self._process_ppt,
            'xlsx': self._process_xlsx,
            'xls': self._process_xls,
            'csv': self._process_csv,
            'txt': self._process_txt,
            'json': self._process_json,
            'rtf': self._process_rtf,
            'odt': self._process_odt,
            'ods': self._process_ods,
            'odp': self._process_odp
        }
        
        # Enhanced processing limits
        self.max_file_size = 100 * 1024 * 1024  # 100MB
        self.max_text_length = 1000000  # 1M characters
        self.max_pages_pdf = 500
        self.max_sheets_excel = 50
    
    def process_file(self, file_path: str, extension: str = None) -> Optional[str]:
        """Process a file and extract its text content with enhanced error handling"""
        try:
            if not os.path.exists(file_path):
                print(f"File not found: {file_path}")
                return None
            
            # Check file size
            file_size = os.path.getsize(file_path)
            if file_size > self.max_file_size:
                print(f"File too large: {file_size} bytes (max: {self.max_file_size})")
                return f"File too large for processing: {file_size / (1024*1024):.1f}MB"
            
            # Determine extension if not provided
            if not extension:
                extension = file_path.split('.')[-1].lower() if '.' in file_path else ''
            
            extension = extension.lower().strip('.')
            
            if extension not in self.supported_extensions:
                print(f"Unsupported file extension: {extension}")
                return f"Unsupported file type: .{extension}"
            
            # Process file based on extension
            processor = self.supported_extensions[extension]
            content = processor(file_path)
            
            if content:
                # Truncate if too long
                if len(content) > self.max_text_length:
                    content = content[:self.max_text_length] + "\n[Content truncated due to length limit]"
                
                print(f"Successfully processed {extension.upper()} file: {os.path.basename(file_path)}")
                return content
            else:
                print(f"No content extracted from: {os.path.basename(file_path)}")
                return f"Could not extract content from {extension.upper()} file"
                
        except Exception as e:
            print(f"Error processing file {file_path}: {e}")
            return f"Error processing file: {str(e)}"
    
    def _process_pdf(self, file_path: str) -> Optional[str]:
        """Extract text from PDF files with enhanced handling"""
        try:
            text_content = []
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = len(pdf_reader.pages)
                
                if num_pages > self.max_pages_pdf:
                    print(f"PDF too long ({num_pages} pages), processing first {self.max_pages_pdf}")
                    num_pages = self.max_pages_pdf
                
                for page_num, page in enumerate(pdf_reader.pages[:num_pages]):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_content.append(f"--- Page {page_num + 1} ---")
                            text_content.append(page_text)
                            text_content.append("")
                    except Exception as e:
                        text_content.append(f"--- Page {page_num + 1} (Error reading) ---")
                        print(f"Error reading PDF page {page_num + 1}: {e}")
                
                if not text_content:
                    # Try alternative PDF processing with pdfplumber if available
                    try:
                        import pdfplumber
                        with pdfplumber.open(file_path) as pdf:
                            for page_num, page in enumerate(pdf.pages[:self.max_pages_pdf]):
                                page_text = page.extract_text()
                                if page_text:
                                    text_content.append(f"--- Page {page_num + 1} ---")
                                    text_content.append(page_text)
                                    text_content.append("")
                    except ImportError:
                        return "PDF contains non-text content or requires advanced processing"
                
                return "\n".join(text_content) if text_content else None
                
        except Exception as e:
            print(f"Error processing PDF: {e}")
            return None
    
    def _process_docx(self, file_path: str) -> Optional[str]:
        """Extract text from DOCX files with enhanced table handling"""
        try:
            doc = docx.Document(file_path)
            text_content = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract tables with improved formatting
            for table_num, table in enumerate(doc.tables, 1):
                text_content.append(f"\n--- Table {table_num} ---")
                for row_num, row in enumerate(table.rows):
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip().replace('\n', ' ').replace('\t', ' ')
                        row_text.append(cell_text)
                    if any(row_text):  # Only add non-empty rows
                        text_content.append(" | ".join(row_text))
                text_content.append("--- End Table ---\n")
            
            return "\n".join(text_content) if text_content else None
            
        except Exception as e:
            print(f"Error processing DOCX: {e}")
            return None
    
    def _process_doc(self, file_path: str) -> Optional[str]:
        """Extract text from DOC files using available tools"""
        try:
            # Try with docx2txt first
            try:
                import docx2txt
                text = docx2txt.process(file_path)
                return text if text.strip() else None
            except ImportError:
                pass
            
            # Try with antiword (if available on system)
            try:
                result = subprocess.run(
                    ['antiword', file_path],
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                if result.returncode == 0:
                    return result.stdout
            except (subprocess.SubprocessError, FileNotFoundError):
                pass
            
            # Try with LibreOffice (if available)
            try:
                temp_dir = tempfile.mkdtemp()
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'txt',
                    '--outdir', temp_dir, file_path
                ], capture_output=True, timeout=60)
                
                if result.returncode == 0:
                    txt_file = os.path.join(temp_dir, os.path.splitext(os.path.basename(file_path))[0] + '.txt')
                    if os.path.exists(txt_file):
                        with open(txt_file, 'r', encoding='utf-8') as f:
                            content = f.read()
                        # Clean up
                        import shutil
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        return content
            except (subprocess.SubprocessError, FileNotFoundError):
                pass
            
            return "DOC file processing requires additional tools (docx2txt, antiword, or LibreOffice)"
            
        except Exception as e:
            print(f"Error processing DOC: {e}")
            return None
    
    def _process_pptx(self, file_path: str) -> Optional[str]:
        """Extract text from PPTX files with enhanced slide handling"""
        try:
            presentation = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []
                slide_text.append(f"--- Slide {slide_num} ---")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Check if it's a title or content
                        if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                            if shape.placeholder_format.type == 1:  # Title
                                slide_text.append(f"TITLE: {shape.text}")
                            else:
                                slide_text.append(shape.text)
                        else:
                            slide_text.append(shape.text)
                
                # Extract table content if present
                for shape in slide.shapes:
                    if shape.has_table:
                        slide_text.append("TABLE:")
                        table = shape.table
                        for row in table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            slide_text.append(" | ".join(row_text))
                
                # Extract notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        slide_text.append(f"NOTES: {notes_text}")
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.extend(slide_text)
                    text_content.append("")
            
            return "\n".join(text_content) if text_content else None
            
        except Exception as e:
            print(f"Error processing PPTX: {e}")
            return None
    
    def _process_ppt(self, file_path: str) -> Optional[str]:
        """Extract text from PPT files using LibreOffice if available"""
        try:
            # Try with LibreOffice
            try:
                temp_dir = tempfile.mkdtemp()
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'txt',
                    '--outdir', temp_dir, file_path
                ], capture_output=True, timeout=60)
                
                if result.returncode == 0:
                    txt_file = os.path.join(temp_dir, os.path.splitext(os.path.basename(file_path))[0] + '.txt')
                    if os.path.exists(txt_file):
                        with open(txt_file, 'r', encoding='utf-8') as f:
                            content = f.read()
                        # Clean up
                        import shutil
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        return content
            except (subprocess.SubprocessError, FileNotFoundError):
                pass
            
            return "PPT file processing requires LibreOffice or conversion to PPTX format"
            
        except Exception as e:
            print(f"Error processing PPT: {e}")
            return None
    
    def _process_xlsx(self, file_path: str) -> Optional[str]:
        """Extract text from XLSX files with enhanced sheet handling"""
        try:
            workbook = load_workbook(file_path, data_only=True)
            text_content = []
            
            sheet_count = 0
            for sheet_name in workbook.sheetnames:
                if sheet_count >= self.max_sheets_excel:
                    text_content.append(f"[Additional {len(workbook.sheetnames) - sheet_count} sheets truncated]")
                    break
                
                sheet = workbook[sheet_name]
                text_content.append(f"--- Sheet: {sheet_name} ---")
                
                # Get data dimensions
                max_row = min(sheet.max_row, 1000)  # Limit rows
                max_col = min(sheet.max_column, 100)  # Limit columns
                
                # Extract data efficiently
                data = []
                for row in sheet.iter_rows(min_row=1, max_row=max_row, min_col=1, max_col=max_col, values_only=True):
                    if any(cell is not None for cell in row):
                        row_data = [str(cell) if cell is not None else "" for cell in row]
                        # Remove trailing empty cells
                        while row_data and not row_data[-1]:
                            row_data.pop()
                        if row_data:
                            data.append(row_data)
                
                if data:
                    # Format as table
                    for row in data[:100]:  # Limit to first 100 rows per sheet
                        text_content.append(" | ".join(row))
                else:
                    text_content.append("[Empty sheet]")
                
                text_content.append("")
                sheet_count += 1
            
            return "\n".join(text_content) if text_content else None
            
        except Exception as e:
            print(f"Error processing XLSX: {e}")
            return None
    
    def _process_xls(self, file_path: str) -> Optional[str]:
        """Extract text from XLS files with enhanced error handling"""
        try:
            # Use pandas to read XLS files
            xl_file = pd.ExcelFile(file_path)
            text_content = []
            
            sheet_count = 0
            for sheet_name in xl_file.sheet_names:
                if sheet_count >= self.max_sheets_excel:
                    text_content.append(f"[Additional {len(xl_file.sheet_names) - sheet_count} sheets truncated]")
                    break
                
                text_content.append(f"--- Sheet: {sheet_name} ---")
                
                try:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    
                    # Convert DataFrame to string representation
                    if not df.empty:
                        # Limit rows and columns for processing
                        limited_df = df.head(100).iloc[:, :20]  # First 100 rows, 20 columns
                        text_content.append(limited_df.to_string(index=False))
                        
                        if len(df) > 100:
                            text_content.append(f"[{len(df) - 100} additional rows not shown]")
                    else:
                        text_content.append("[Empty sheet]")
                        
                except Exception as e:
                    text_content.append(f"[Error reading sheet: {e}]")
                
                text_content.append("")
                sheet_count += 1
            
            return "\n".join(text_content) if text_content else None
            
        except Exception as e:
            print(f"Error processing XLS: {e}")
            return None
    
    def _process_csv(self, file_path: str) -> Optional[str]:
        """Extract text from CSV files with enhanced encoding detection"""
        try:
            # Detect encoding
            encoding = 'utf-8'
            try:
                with open(file_path, 'rb') as f:
                    raw_data = f.read(10000)  # Read first 10KB
                    detected = chardet.detect(raw_data)
                    if detected['encoding'] and detected['confidence'] > 0.7:
                        encoding = detected['encoding']
            except:
                pass
            
            # Try reading with detected encoding
            try:
                df = pd.read_csv(file_path, encoding=encoding)
            except UnicodeDecodeError:
                # Fallback encodings
                for fallback_encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                    try:
                        df = pd.read_csv(file_path, encoding=fallback_encoding)
                        break
                    except UnicodeDecodeError:
                        continue
                else:
                    return "Could not decode CSV file with any encoding"
            
            text_content = []
            text_content.append("--- CSV Data ---")
            text_content.append(f"Columns ({len(df.columns)}): {', '.join(df.columns.astype(str).tolist())}")
            text_content.append(f"Total rows: {len(df)}")
            text_content.append("")
            
            # Show first 100 rows and limit columns if too many
            limited_df = df.head(100)
            if len(df.columns) > 20:
                limited_df = limited_df.iloc[:, :20]
                text_content.append(f"[Showing first 20 of {len(df.columns)} columns]")
            
            text_content.append(limited_df.to_string(index=False))
            
            if len(df) > 100:
                text_content.append(f"\n[{len(df) - 100} additional rows not shown]")
            
            return "\n".join(text_content)
            
        except Exception as e:
            print(f"Error processing CSV: {e}")
            return None
    
    def _process_txt(self, file_path: str) -> Optional[str]:
        """Extract text from TXT files with encoding detection"""
        try:
            # Detect encoding
            encoding = 'utf-8'
            try:
                with open(file_path, 'rb') as f:
                    raw_data = f.read()
                    detected = chardet.detect(raw_data)
                    if detected['encoding'] and detected['confidence'] > 0.7:
                        encoding = detected['encoding']
            except:
                pass
            
            # Try reading with detected encoding first
            encodings_to_try = [encoding, 'utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for enc in encodings_to_try:
                try:
                    with open(file_path, 'r', encoding=enc) as file:
                        content = file.read()
                        return content if content.strip() else None
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, try with error handling
            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                return file.read()
                
        except Exception as e:
            print(f"Error processing TXT: {e}")
            return None
    
    def _process_json(self, file_path: str) -> Optional[str]:
        """Extract text from JSON files with pretty formatting"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
            
            # Convert JSON to readable text format
            if isinstance(data, dict):
                text_content = ["--- JSON Object ---"]
                text_content.append(json.dumps(data, indent=2, ensure_ascii=False)[:50000])  # Limit size
            elif isinstance(data, list):
                text_content = ["--- JSON Array ---"]
                text_content.append(f"Array with {len(data)} items:")
                sample_items = min(10, len(data))
                text_content.append(json.dumps(data[:sample_items], indent=2, ensure_ascii=False))
                if len(data) > sample_items:
                    text_content.append(f"... and {len(data) - sample_items} more items")
            else:
                text_content = [str(data)]
            
            return "\n".join(text_content)
            
        except Exception as e:
            print(f"Error processing JSON: {e}")
            return None
    
    def _process_rtf(self, file_path: str) -> Optional[str]:
        """Extract text from RTF files"""
        try:
            # Try with striprtf if available
            try:
                from striprtf.striprtf import rtf_to_text
                with open(file_path, 'r', encoding='utf-8') as file:
                    rtf_content = file.read()
                return rtf_to_text(rtf_content)
            except ImportError:
                pass
            
            # Try with LibreOffice
            try:
                temp_dir = tempfile.mkdtemp()
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'txt',
                    '--outdir', temp_dir, file_path
                ], capture_output=True, timeout=60)
                
                if result.returncode == 0:
                    txt_file = os.path.join(temp_dir, os.path.splitext(os.path.basename(file_path))[0] + '.txt')
                    if os.path.exists(txt_file):
                        with open(txt_file, 'r', encoding='utf-8') as f:
                            content = f.read()
                        # Clean up
                        import shutil
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        return content
            except (subprocess.SubprocessError, FileNotFoundError):
                pass
            
            return "RTF file processing requires striprtf package or LibreOffice"
            
        except Exception as e:
            print(f"Error processing RTF: {e}")
            return None
    
    def _process_odt(self, file_path: str) -> Optional[str]:
        """Extract text from ODT files using LibreOffice or zip extraction"""
        try:
            # Try with LibreOffice first
            try:
                temp_dir = tempfile.mkdtemp()
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'txt',
                    '--outdir', temp_dir, file_path
                ], capture_output=True, timeout=60)
                
                if result.returncode == 0:
                    txt_file = os.path.join(temp_dir, os.path.splitext(os.path.basename(file_path))[0] + '.txt')
                    if os.path.exists(txt_file):
                        with open(txt_file, 'r', encoding='utf-8') as f:
                            content = f.read()
                        # Clean up
                        import shutil
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        return content
            except (subprocess.SubprocessError, FileNotFoundError):
                pass
            
            # Try manual XML extraction
            try:
                with zipfile.ZipFile(file_path, 'r') as zip_file:
                    if 'content.xml' in zip_file.namelist():
                        content_xml = zip_file.read('content.xml').decode('utf-8')
                        # Basic XML tag removal (not perfect but functional)
                        import re
                        text = re.sub(r'<[^>]+>', ' ', content_xml)
                        text = re.sub(r'\s+', ' ', text)
                        return text.strip() if text.strip() else None
            except Exception:
                pass
            
            return "ODT file processing requires LibreOffice"
            
        except Exception as e:
            print(f"Error processing ODT: {e}")
            return None
    
    def _process_ods(self, file_path: str) -> Optional[str]:
        """Extract text from ODS files"""
        try:
            # Try with LibreOffice
            try:
                temp_dir = tempfile.mkdtemp()
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'csv',
                    '--outdir', temp_dir, file_path
                ], capture_output=True, timeout=60)
                
                if result.returncode == 0:
                    csv_file = os.path.join(temp_dir, os.path.splitext(os.path.basename(file_path))[0] + '.csv')
                    if os.path.exists(csv_file):
                        content = self._process_csv(csv_file)
                        # Clean up
                        import shutil
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        return content
            except (subprocess.SubprocessError, FileNotFoundError):
                pass
            
            return "ODS file processing requires LibreOffice"
        except Exception as e:
            print(f"Error processing ODS: {e}")
            return None
    
    def _process_odp(self, file_path: str) -> Optional[str]:
        """Extract text from ODP files"""
        try:
            # Try with LibreOffice
            try:
                temp_dir = tempfile.mkdtemp()
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'txt',
                    '--outdir', temp_dir, file_path
                ], capture_output=True, timeout=60)
                
                if result.returncode == 0:
                    txt_file = os.path.join(temp_dir, os.path.splitext(os.path.basename(file_path))[0] + '.txt')
                    if os.path.exists(txt_file):
                        with open(txt_file, 'r', encoding='utf-8') as f:
                            content = f.read()
                        # Clean up
                        import shutil
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        return content
            except (subprocess.SubprocessError, FileNotFoundError):
                pass
            
            return "ODP file processing requires LibreOffice"
        except Exception as e:
            print(f"Error processing ODP: {e}")
            return None
    
    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """Get comprehensive information about a file"""
        try:
            stat = os.stat(file_path)
            extension = file_path.split('.')[-1].lower() if '.' in file_path else ''
            
            return {
                'filename': os.path.basename(file_path),
                'size': stat.st_size,
                'size_mb': round(stat.st_size / (1024 * 1024), 2),
                'size_human': self._format_file_size(stat.st_size),
                'extension': extension,
                'supported': extension in self.supported_extensions,
                'modified': stat.st_mtime,
                'type': self._get_file_type(extension),
                'processing_complexity': self._get_processing_complexity(extension, stat.st_size)
            }
        except Exception as e:
            return {
                'filename': os.path.basename(file_path) if file_path else 'unknown',
                'error': str(e),
                'supported': False,
                'type': 'unknown'
            }
    
    def _format_file_size(self, size_bytes: int) -> str:
        """Format file size in human readable format"""
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_bytes < 1024.0:
                return f"{size_bytes:.1f} {unit}"
            size_bytes /= 1024.0
        return f"{size_bytes:.1f} TB"
    
    def _get_file_type(self, extension: str) -> str:
        """Get file type category"""
        document_types = {'pdf', 'docx', 'doc', 'txt', 'rtf', 'odt'}
        spreadsheet_types = {'xlsx', 'xls', 'csv', 'ods'}
        presentation_types = {'pptx', 'ppt', 'odp'}
        data_types = {'json', 'xml'}
        
        if extension in document_types:
            return 'document'
        elif extension in spreadsheet_types:
            return 'spreadsheet'
        elif extension in presentation_types:
            return 'presentation'
        elif extension in data_types:
            return 'data'
        else:
            return 'unknown'
    
    def _get_processing_complexity(self, extension: str, file_size: int) -> str:
        """Estimate processing complexity"""
        if extension in ['txt', 'csv', 'json']:
            return 'low'
        elif extension in ['docx', 'xlsx', 'pptx'] and file_size < 10 * 1024 * 1024:
            return 'medium'
        elif extension in ['pdf', 'doc', 'xls', 'ppt'] or file_size > 10 * 1024 * 1024:
            return 'high'
        else:
            return 'medium'
    
    def batch_process_files(self, file_paths: List[str]) -> Dict[str, Any]:
        """Process multiple files and return comprehensive results"""
        results = {
            'successful': [],
            'failed': [],
            'combined_content': [],
            'total_files': len(file_paths),
            'total_size': 0,
            'processing_time': 0,
            'file_types': {}
        }
        
        import time
        start_time = time.time()
        
        for file_path in file_paths:
            try:
                file_info = self.get_file_info(file_path)
                results['total_size'] += file_info.get('size', 0)
                
                file_type = file_info.get('type', 'unknown')
                results['file_types'][file_type] = results['file_types'].get(file_type, 0) + 1
                
                if file_info.get('supported', False):
                    content = self.process_file(file_path)
                    if content:
                        results['successful'].append({
                            'filename': file_info['filename'],
                            'content': content,
                            'size_mb': file_info['size_mb'],
                            'type': file_type,
                            'complexity': file_info.get('processing_complexity', 'unknown')
                        })
                        results['combined_content'].append(f"=== {file_info['filename']} ===")
                        results['combined_content'].append(content)
                        results['combined_content'].append("")
                    else:
                        results['failed'].append({
                            'filename': file_info['filename'],
                            'reason': 'No content extracted',
                            'type': file_type
                        })
                else:
                    results['failed'].append({
                        'filename': file_info['filename'],
                        'reason': 'Unsupported file type',
                        'type': file_type
                    })
                    
            except Exception as e:
                results['failed'].append({
                    'filename': os.path.basename(file_path) if file_path else 'unknown',
                    'reason': str(e),
                    'type': 'unknown'
                })
        
        results['processing_time'] = time.time() - start_time
        results['combined_text'] = "\n".join(results['combined_content'])
        results['success_rate'] = len(results['successful']) / len(file_paths) if file_paths else 0
        results['total_size_mb'] = results['total_size'] / (1024 * 1024)
        
        return results
    
    def validate_file(self, file_path: str) -> Dict[str, Any]:
        """Validate a file before processing"""
        validation_result = {
            'valid': False,
            'errors': [],
            'warnings': [],
            'info': {}
        }
        
        try:
            if not os.path.exists(file_path):
                validation_result['errors'].append("File does not exist")
                return validation_result
            
            file_info = self.get_file_info(file_path)
            validation_result['info'] = file_info
            
            # Check file size
            if file_info['size'] > self.max_file_size:
                validation_result['errors'].append(f"File too large: {file_info['size_human']} (max: {self._format_file_size(self.max_file_size)})")
            
            if file_info['size'] == 0:
                validation_result['errors'].append("File is empty")
            
            # Check file type
            if not file_info['supported']:
                validation_result['errors'].append(f"Unsupported file type: .{file_info['extension']}")
            
            # Check processing complexity
            complexity = file_info.get('processing_complexity', 'unknown')
            if complexity == 'high':
                validation_result['warnings'].append("File may require significant processing time")
            
            # Additional format-specific validation
            extension = file_info['extension']
            if extension == 'pdf' and file_info['size'] > 50 * 1024 * 1024:
                validation_result['warnings'].append("Large PDF files may have incomplete text extraction")
            
            validation_result['valid'] = len(validation_result['errors']) == 0
            
        except Exception as e:
            validation_result['errors'].append(f"Validation error: {str(e)}")
        
        return validation_result